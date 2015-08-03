from httplib2 import Http
from lxml.html import document_fromstring
import sys
from StringIO import StringIO
from argparse import ArgumentParser
from docx import Document
from docx.document import Document as _Document
from pptx import util
from pptx import Presentation
from pptx.api import Presentation as _Presentation
from pptx.parts import slide as _Slide
from pptx import enum

if __name__ == "__main__":
  argparser = ArgumentParser(description="convert docx text into pptx slide notes")
  argparser.add_argument(
    '-f', '--input',
    dest='input',
  )
  argparser.add_argument(
    '-o', '--output',
    dest='output',
  )

  args = argparser.parse_args()

  fdIn = sys.stdin
  if (args.input != None):
    fdIn = open(args.input)

  fdOut = sys.stdout
  if (args.output != None):
    fdOut = open(args.output, 'w')

  html = fdIn.read()
  dom = document_fromstring(html)
  fdIn.close()

  pptx = Presentation()
  layout = pptx.slide_layouts[2]


  rows = dom.xpath("descendant::tr")
  for row in rows:
      cells = row.xpath("td")

      # TODO: replace this HACK with some nice UI?
      if len(cells) == 3:
          script_cell = cells[1]
          visuals_cell = cells[2]
      elif len(cells) == 2:
          script_cell = cells[0]
          visuals_cell = cells[1]
      else:
          continue

      slide = pptx.slides.add_slide(layout)

      text = script_cell.xpath('normalize-space(.)')
      if text and len(text) > 0:
          slide.notes.add_note(text)

      text = visuals_cell.xpath('normalize-space(.)')
      if text and len(text) > 0:
          tbox = slide.shapes.add_textbox(util.Inches(0), util.Inches(0), pptx.slide_width, util.Inches(2))
          tbox.auto_size = enum.text.MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
          tbox.text = text
      images = visuals_cell.xpath("descendant::img")
      for img in images:
          src = img.attrib.get('src')
          r, bits = Http().request(src)
          fbits = StringIO(bits)
          ibox = slide.shapes.add_picture(fbits, 0, 0, pptx.slide_width, pptx.slide_height)


  pptx.save(fdOut)

  fdOut.close()

